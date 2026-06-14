import os
import io
import re
from typing import Dict, Any
from app.services.llm_service import llm_service
from datetime import datetime

class AdvancedReportService:
    @staticmethod
    def generate_report_content(data: Dict[str, Any], report_type: str) -> str:
        """Uses LLM to generate the text content tailored to the specific report type."""
        
        prompts = {
            "Executive Summary": "You are a Data Science VP writing a concise Executive Summary for the C-Suite. Focus only on high-level strategic takeaways, bottom-line impact, and top 3 recommendations.",
            "Board Report": "You are preparing a formal Board Report. Include an executive summary, detailed financial/operational findings, risk analysis based on anomalies, and strategic next steps. Use a very formal, professional tone.",
            "KPI Review": "You are an Analytics Manager writing a KPI Review. Focus heavily on the metrics. Highlight period-over-period trends, outliers, and operational bottlenecks. Use bullet points extensively.",
            "Market Analysis": "You are a Market Researcher. Synthesize this data into a Market Analysis report. Focus on overall trends, comparative performance, forecasting future market states, and competitive positioning."
        }
        
        system_prompt = f"""
        {prompts.get(report_type, prompts['Executive Summary'])}
        
        Raw Analytics Data:
        {data}
        
        Format the output purely in Markdown. Use headers (##), bold text (**), and bullet points (-). 
        Do not include conversational filler like 'Here is your report'.
        """
        
        try:
            return llm_service.llm.invoke(system_prompt).content
        except Exception:
            return f"## {report_type}\n*Generated {datetime.now().strftime('%Y-%m-%d')}*\n\n### Summary\nDue to missing LLM keys, this is a placeholder report. The automated systems have successfully compiled the raw data.\n\n### Data Snapshot\n{str(data)[:500]}..."

    @staticmethod
    def _export_to_docx(markdown_text: str) -> io.BytesIO:
        from docx import Document
        from docx.shared import Pt
        
        doc = Document()
        doc.add_heading('AntiGravity Platform Report', 0)
        
        # Simple Markdown parser for Word
        lines = markdown_text.split('\n')
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            if line.startswith('# '):
                doc.add_heading(line[2:], level=1)
            elif line.startswith('## '):
                doc.add_heading(line[3:], level=2)
            elif line.startswith('### '):
                doc.add_heading(line[4:], level=3)
            elif line.startswith('- ') or line.startswith('* '):
                p = doc.add_paragraph(style='List Bullet')
                # Handle bold text roughly
                parts = re.split(r'\*\*(.*?)\*\*', line[2:])
                for i, part in enumerate(parts):
                    if i % 2 == 1:
                        p.add_run(part).bold = True
                    else:
                        p.add_run(part)
            else:
                p = doc.add_paragraph()
                parts = re.split(r'\*\*(.*?)\*\*', line)
                for i, part in enumerate(parts):
                    if i % 2 == 1:
                        p.add_run(part).bold = True
                    else:
                        p.add_run(part)
        
        file_stream = io.BytesIO()
        doc.save(file_stream)
        file_stream.seek(0)
        return file_stream

    @staticmethod
    def _export_to_pdf(markdown_text: str) -> io.BytesIO:
        from fpdf import FPDF
        
        class ReportPDF(FPDF):
            def header(self):
                self.set_font('Arial', 'B', 12)
                self.cell(0, 10, 'AntiGravity Platform Report', 0, 1, 'C')
                self.line(10, 20, 200, 20)
                self.ln(10)
                
            def footer(self):
                self.set_y(-15)
                self.set_font('Arial', 'I', 8)
                self.cell(0, 10, f'Page {self.page_no()}', 0, 0, 'C')
        
        pdf = ReportPDF()
        pdf.add_page()
        pdf.set_auto_page_break(auto=True, margin=15)
        
        lines = markdown_text.split('\n')
        for line in lines:
            line = line.strip()
            if not line:
                pdf.ln(5)
                continue
                
            # Very basic markdown handling for PDF
            line = line.replace('**', '') # fpdf2 handles bold via set_font, stripping for simplicity in baseline
            
            if line.startswith('# '):
                pdf.set_font('Arial', 'B', 18)
                pdf.multi_cell(0, 10, line[2:])
                pdf.ln(2)
            elif line.startswith('## '):
                pdf.set_font('Arial', 'B', 14)
                pdf.multi_cell(0, 10, line[3:])
                pdf.ln(2)
            elif line.startswith('### '):
                pdf.set_font('Arial', 'B', 12)
                pdf.multi_cell(0, 10, line[4:])
                pdf.ln(2)
            elif line.startswith('- ') or line.startswith('* '):
                pdf.set_font('Arial', '', 11)
                pdf.multi_cell(0, 7, f"    {line}")
            else:
                pdf.set_font('Arial', '', 11)
                pdf.multi_cell(0, 7, line)
                
        # fpdf2 outputs string natively, we encode it to bytes
        file_stream = io.BytesIO(pdf.output(dest='S').encode('latin-1', 'replace'))
        file_stream.seek(0)
        return file_stream

    @staticmethod
    def _export_to_pptx(markdown_text: str) -> io.BytesIO:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        
        prs = Presentation()
        
        # Title Slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "AntiGravity Platform Report"
        subtitle.text = f"Automated AI Generation\n{datetime.now().strftime('%Y-%m-%d')}"
        
        # Split markdown into sections by H2
        sections = re.split(r'\n## ', markdown_text)
        
        for section in sections:
            if not section.strip() or section.startswith('# '):
                continue
                
            lines = section.split('\n')
            slide_title = lines[0].strip().replace('*', '')
            
            content_layout = prs.slide_layouts[1] # Title and Content
            slide = prs.slides.add_slide(content_layout)
            title_shape = slide.shapes.title
            body_shape = slide.placeholders[1]
            
            title_shape.text = slide_title
            tf = body_shape.text_frame
            
            bullet_points_added = 0
            for line in lines[1:]:
                line = line.strip().replace('**', '')
                if not line: continue
                
                # Limit to 5 bullets per slide to avoid overflow
                if bullet_points_added > 5:
                    break
                    
                if line.startswith('- ') or line.startswith('* '):
                    p = tf.add_paragraph()
                    p.text = line[2:]
                    p.level = 0
                    bullet_points_added += 1
                elif not line.startswith('#'):
                    p = tf.add_paragraph()
                    p.text = line
                    p.level = 0
                    bullet_points_added += 1
        
        file_stream = io.BytesIO()
        prs.save(file_stream)
        file_stream.seek(0)
        return file_stream

    @staticmethod
    def generate_document(data: Dict[str, Any], report_type: str, format_type: str) -> io.BytesIO:
        # 1. Generate text using LLM
        markdown_text = AdvancedReportService.generate_report_content(data, report_type)
        
        # 2. Convert to requested format
        format_type = format_type.upper()
        if format_type == "DOCX":
            return AdvancedReportService._export_to_docx(markdown_text)
        elif format_type == "PDF":
            return AdvancedReportService._export_to_pdf(markdown_text)
        elif format_type == "PPTX":
            return AdvancedReportService._export_to_pptx(markdown_text)
        else:
            # Fallback to pure text/markdown
            stream = io.BytesIO(markdown_text.encode('utf-8'))
            stream.seek(0)
            return stream
